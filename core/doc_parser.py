import os
from pptx import Presentation
from pypdf import PdfReader
from docx import Document

def extract_text_from_pptx(file_stream):
    text_content = []
    try:
        file_stream.seek(0)
        prs = Presentation(file_stream)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text_content.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
    except Exception as e:
        return f"[PPTX解析错误: {str(e)}]"
    return "\n\n".join(text_content)

def extract_text_from_pdf(file_stream):
    text_content = []
    try:
        file_stream.seek(0)
        reader = PdfReader(file_stream)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_content.append(f"[Page {i+1}]\n{page_text}")
    except Exception as e:
        return f"[PDF解析错误: {str(e)}]"
    return "\n\n".join(text_content)

def extract_text_from_docx(file_stream):
    text_content = []
    try:
        file_stream.seek(0)
        doc = Document(file_stream)
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
    except Exception as e:
        return f"[DOCX解析错误: {str(e)}]"
    return "\n".join(text_content)

def extract_text_from_plain(file_stream, encoding='utf-8'):
    try:
        file_stream.seek(0)
        return file_stream.read().decode(encoding)
    except Exception:
        try:
            file_stream.seek(0)
            return file_stream.read().decode('gbk')
        except Exception as e:
            return f"[文本解析错误: {str(e)}]"

def parse_reference_files(uploaded_files):
    all_reference_content = ""
    
    if not uploaded_files:
        return ""

    for uploaded_file in uploaded_files:
        filename = uploaded_file.name
        file_ext = os.path.splitext(filename)[1].lower()
        content = ""
        
        try:
            if file_ext == '.pptx':
                content = extract_text_from_pptx(uploaded_file)
            elif file_ext == '.pdf':
                content = extract_text_from_pdf(uploaded_file)
            elif file_ext == '.docx':
                content = extract_text_from_docx(uploaded_file)
            elif file_ext in ['.txt', '.md']:
                content = extract_text_from_plain(uploaded_file)
            else:
                continue
                
            if content.strip():
                all_reference_content += f"\n\n=== 参考资料: {filename} ===\n{content}\n"
                
        except Exception as e:
            all_reference_content += f"\n\n=== 参考资料: {filename} (读取失败) ===\n错误信息: {str(e)}\n"

    return all_reference_content